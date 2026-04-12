"""파일 업로드 / 다운로드 라우트"""

import uuid
from pathlib import Path

from fastapi import APIRouter, HTTPException, UploadFile, File, Header
from fastapi.responses import FileResponse

import structlog

from core.config import settings

logger = structlog.get_logger(__name__)
router = APIRouter(prefix="/api", tags=["files"])


def _extract_text(path: Path, filename: str, content: bytes) -> str:
    """파일 내용에서 텍스트 추출 (형식별 처리)"""
    ext = Path(filename).suffix.lower()
    try:
        if ext == ".docx":
            from docx import Document  # type: ignore
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        elif ext == ".xlsx":
            import openpyxl  # type: ignore
            wb = openpyxl.load_workbook(path, data_only=True)
            rows: list[str] = []
            for sheet in wb.worksheets:
                rows.append(f"[시트: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        rows.append(row_text)
            return "\n".join(rows)

        elif ext == ".pptx":
            from pptx import Presentation  # type: ignore
            prs = Presentation(path)
            texts: list[str] = []
            for i, slide in enumerate(prs.slides, 1):
                texts.append(f"[슬라이드 {i}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts)

        elif ext == ".pdf":
            try:
                from pypdf import PdfReader  # type: ignore
                reader = PdfReader(path)
                pages = [page.extract_text() or "" for page in reader.pages]
                return "\n\n".join(p for p in pages if p.strip())
            except ImportError:
                return f"[PDF 업로드됨: {filename} — pypdf 미설치]"

        elif ext in (".txt", ".md", ".csv"):
            return content.decode("utf-8", errors="ignore")

        else:
            return f"[파일 업로드됨: {filename}]"

    except Exception as exc:
        logger.warning("text_extraction_failed", filename=filename, error=str(exc))
        return f"[파일 업로드됨: {filename} — 텍스트 추출 불가]"


# 확장자 → MIME 타입 매핑
MIME_TYPES = {
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pdf": "application/pdf",
}


@router.post("/upload")
async def upload_file(
    file: UploadFile = File(...),
    x_user_id: str | None = Header(default=None),
):
    """파일 업로드 및 RAG 파이프라인 실행

    1. 파일 저장 및 텍스트 추출
    2. 텍스트를 청크로 분할 (800자 단위)
    3. 각 청크를 OpenAI 임베딩으로 벡터화
    4. Supabase document_chunks 테이블에 저장
    5. 응답에 upload_id, chunk_count 포함
    """
    upload_id = str(uuid.uuid4())
    upload_dir = settings.upload_dir
    upload_dir.mkdir(parents=True, exist_ok=True)

    safe_name = Path(file.filename or "upload").name
    save_path = upload_dir / f"{upload_id}_{safe_name}"

    content = await file.read()
    save_path.write_bytes(content)

    # 텍스트 추출 (PDF 포함)
    extracted = _extract_text(save_path, safe_name, content)

    logger.info(
        "file_uploaded",
        upload_id=upload_id,
        filename=safe_name,
        size=len(content),
        text_length=len(extracted),
    )

    # ── RAG 파이프라인 (비동기 — 실패해도 업로드는 성공 처리) ──
    chunk_count = 0
    rag_enabled = False

    try:
        from rag.chunker import chunk_text
        from rag.embedder import embed_texts
        from rag.store import save_chunks

        chunks = chunk_text(extracted, chunk_size=800, overlap=100)
        chunk_count = len(chunks)

        if chunks:
            vectors = await embed_texts(chunks)
            await save_chunks(
                upload_id=upload_id,
                user_id=x_user_id,
                chunks=chunks,
                vectors=vectors,
            )
            rag_enabled = True
            logger.info("rag_pipeline_done", upload_id=upload_id, chunk_count=chunk_count)

    except Exception as exc:
        # RAG 실패 시 fallback: 기존 방식으로 4000자 직접 제공
        logger.warning("rag_pipeline_failed", upload_id=upload_id, error=str(exc))

    # fallback: RAG 실패 시 기존 방식 호환
    extracted_trimmed = extracted[:4000] if not rag_enabled else extracted[:500]

    return {
        "upload_id": upload_id,
        "filename": safe_name,
        "size_bytes": len(content),
        "extracted_text": extracted_trimmed,
        "chunk_count": chunk_count,
        "rag_enabled": rag_enabled,
    }


@router.get("/files/{file_id}")
async def download_file(file_id: str):
    """생성된 파일 다운로드

    UUID 기반 파일 경로로 접근합니다.
    파일은 생성 후 1시간 뒤 자동 삭제됩니다.
    """
    for ext in MIME_TYPES:
        file_path = settings.output_dir / f"{file_id}{ext}"
        if file_path.exists():
            media_type = MIME_TYPES[ext]

            logger.info(
                "file_download",
                file_id=file_id,
                ext=ext,
                size=file_path.stat().st_size,
            )

            return FileResponse(
                path=str(file_path),
                media_type=media_type,
                filename=f"docflow_{file_id[:8]}{ext}",
                headers={
                    "Content-Disposition": f'attachment; filename="docflow_{file_id[:8]}{ext}"',
                },
            )

    raise HTTPException(
        status_code=404,
        detail="파일을 찾을 수 없습니다. 만료되었거나 존재하지 않는 파일입니다.",
    )
