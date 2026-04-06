"""PPT 생성 엔진 — python-pptx 래퍼"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import structlog

logger = structlog.get_logger(__name__)

# ── 테마 색상 정의 ──
THEMES: dict[str, dict[str, RGBColor]] = {
    "modern": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title": RGBColor(0x1A, 0x1A, 0x2E),
        "accent": RGBColor(0x4A, 0x90, 0xD9),
        "text": RGBColor(0x33, 0x33, 0x33),
        "subtitle": RGBColor(0x66, 0x66, 0x66),
    },
    "classic": {
        "bg": RGBColor(0xFD, 0xFD, 0xFD),
        "title": RGBColor(0x2C, 0x3E, 0x50),
        "accent": RGBColor(0xE7, 0x4C, 0x3C),
        "text": RGBColor(0x2C, 0x3E, 0x50),
        "subtitle": RGBColor(0x7F, 0x8C, 0x8D),
    },
    "minimal": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title": RGBColor(0x00, 0x00, 0x00),
        "accent": RGBColor(0x00, 0x00, 0x00),
        "text": RGBColor(0x44, 0x44, 0x44),
        "subtitle": RGBColor(0x88, 0x88, 0x88),
    },
    "dark": {
        "bg": RGBColor(0x1A, 0x1A, 0x2E),
        "title": RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0x4A, 0x90, 0xD9),
        "text": RGBColor(0xE0, 0xE0, 0xE0),
        "subtitle": RGBColor(0xAA, 0xAA, 0xAA),
    },
}

# ── 한국어 폰트 ──
FONT_NAME = "맑은 고딕"
FONT_NAME_EN = "Pretendard"


class PptEngine:
    """python-pptx 기반 PPT 생성 엔진"""

    def __init__(self, theme: str = "modern"):
        self.prs = Presentation()
        self.theme = THEMES.get(theme, THEMES["modern"])

        # 16:9 와이드스크린
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

    def create(
        self, title: str, slides: list[dict], output_path: Path
    ) -> dict:
        """PPT 파일 생성

        Args:
            title: 프레젠테이션 전체 제목
            slides: 슬라이드 데이터 목록
            output_path: 저장 경로

        Returns:
            dict: 생성 결과 정보
        """
        # 타이틀 슬라이드
        self._add_title_slide(title)

        # 콘텐츠 슬라이드
        for slide_data in slides:
            layout = slide_data.get("layout", "content")
            if layout == "two_column":
                self._add_two_column_slide(slide_data)
            elif layout == "title_only":
                self._add_title_only_slide(slide_data)
            else:
                self._add_content_slide(slide_data)

        self.prs.save(str(output_path))
        file_size = output_path.stat().st_size

        logger.info(
            "ppt_created",
            title=title,
            slides=len(slides) + 1,
            size_bytes=file_size,
            path=str(output_path),
        )

        return {
            "slide_count": len(slides) + 1,
            "size_bytes": file_size,
        }

    def _add_title_slide(self, title: str) -> None:
        """타이틀 슬라이드 추가"""
        layout = self.prs.slide_layouts[0]  # Title Slide
        slide = self.prs.slides.add_slide(layout)

        # 제목
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(40)
                run.font.bold = True
                run.font.color.rgb = self.theme["title"]
                run.font.name = FONT_NAME

        # 부제목 (있으면)
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = "DocFlow AI로 자동 생성"
            for para in subtitle.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(18)
                    run.font.color.rgb = self.theme["subtitle"]
                    run.font.name = FONT_NAME

    def _add_content_slide(self, data: dict) -> None:
        """일반 콘텐츠 슬라이드 (제목 + 불릿 포인트)"""
        layout = self.prs.slide_layouts[1]  # Title and Content
        slide = self.prs.slides.add_slide(layout)

        # 제목
        title_shape = slide.shapes.title
        title_shape.text = data.get("title", "")
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(32)
                run.font.bold = True
                run.font.color.rgb = self.theme["title"]
                run.font.name = FONT_NAME

        # 내용 불릿 포인트
        content_items = data.get("content", [])
        if content_items and len(slide.placeholders) > 1:
            body = slide.placeholders[1].text_frame
            body.clear()

            for i, item in enumerate(content_items):
                para = body.paragraphs[0] if i == 0 else body.add_paragraph()
                para.text = item
                para.level = 0
                para.space_after = Pt(8)
                for run in para.runs:
                    run.font.size = Pt(20)
                    run.font.color.rgb = self.theme["text"]
                    run.font.name = FONT_NAME

        # 발표자 노트
        notes = data.get("notes")
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def _add_two_column_slide(self, data: dict) -> None:
        """2열 레이아웃 슬라이드"""
        layout = self.prs.slide_layouts[3]  # Two Content (보통 index 3)
        slide = self.prs.slides.add_slide(layout)

        # 제목
        title_shape = slide.shapes.title
        title_shape.text = data.get("title", "")
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(32)
                run.font.bold = True
                run.font.color.rgb = self.theme["title"]
                run.font.name = FONT_NAME

        # 좌측 / 우측 content
        content = data.get("content", [])
        mid = len(content) // 2

        for idx, placeholder_idx in enumerate([1, 2]):
            if placeholder_idx < len(slide.placeholders):
                items = content[:mid] if idx == 0 else content[mid:]
                body = slide.placeholders[placeholder_idx].text_frame
                body.clear()
                for i, item in enumerate(items):
                    para = body.paragraphs[0] if i == 0 else body.add_paragraph()
                    para.text = item
                    for run in para.runs:
                        run.font.size = Pt(18)
                        run.font.color.rgb = self.theme["text"]
                        run.font.name = FONT_NAME

    def _add_title_only_slide(self, data: dict) -> None:
        """제목만 있는 슬라이드"""
        layout = self.prs.slide_layouts[5]  # Blank
        slide = self.prs.slides.add_slide(layout)

        from pptx.util import Emu

        # 수동으로 제목 텍스트박스 추가
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(11.33)
        height = Inches(2)

        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True

        para = tf.paragraphs[0]
        para.text = data.get("title", "")
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.color.rgb = self.theme["title"]
            run.font.name = FONT_NAME
