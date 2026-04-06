from pathlib import Path
import io
import asyncio
import os

async def export_to_pdf(slides_dir: Path, output_path: Path) -> None:
    """HTML 슬라이드 목록을 단일 PDF로 병합"""
    try:
        from playwright.async_api import async_playwright
        from pypdf import PdfWriter, PdfReader
    except ImportError as e:
        raise RuntimeError(f"PDF 내보내기에 필요한 패키지가 없습니다: {e}") from e

    async with async_playwright() as p:
        browser = await p.chromium.launch(args=["--no-sandbox"])
        page = await browser.new_page(
            viewport={"width": 1280, "height": 720}
        )

        writer = PdfWriter()
        slide_files = sorted(slides_dir.glob("slide-*.html"))

        for slide_file in slide_files:
            # 절대 경로를 file:// URI로 변환 (Windows 대응)
            import urllib.parse
            path_uri = urllib.parse.quote(str(slide_file.absolute()).replace('\\', '/'))
            # 드라이브 문자(ex. C:/)로 시작할 경우, 맨 앞에 / 추가
            if ':' in path_uri.split('/')[0]:
                path_uri = '/' + path_uri
                
            await page.goto(
                f"file://{path_uri}",
                wait_until="networkidle",
                timeout=15000,
            )
            await page.wait_for_timeout(300)

            pdf_bytes = await page.pdf(
                width="1280px",
                height="720px",
                print_background=True,
                margin={"top": "0", "right": "0", "bottom": "0", "left": "0"},
            )

            reader = PdfReader(io.BytesIO(pdf_bytes))
            for pdf_page in reader.pages:
                writer.add_page(pdf_page)

        await browser.close()

    output_path.parent.mkdir(parents=True, exist_ok=True)
    with open(output_path, "wb") as f:
        writer.write(f)


async def export_to_pptx(slides_dir: Path, output_path: Path) -> None:
    """HTML 슬라이드를 스크린샷으로 캡처해 PPTX 생성"""
    try:
        from playwright.async_api import async_playwright
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as e:
        raise RuntimeError(f"PPTX 내보내기에 필요한 패키지가 없습니다: {e}") from e

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    async with async_playwright() as p:
        browser = await p.chromium.launch(args=["--no-sandbox"])
        page = await browser.new_page(
            viewport={"width": 1280, "height": 720}
        )

        slide_files = sorted(slides_dir.glob("slide-*.html"))

        for slide_file in slide_files:
            import urllib.parse
            path_uri = urllib.parse.quote(str(slide_file.absolute()).replace('\\', '/'))
            if ':' in path_uri.split('/')[0]:
                path_uri = '/' + path_uri

            await page.goto(
                f"file://{path_uri}",
                wait_until="networkidle",
                timeout=15000,
            )
            await page.wait_for_timeout(400)

            screenshot = await page.screenshot(
                full_page=False,
                type="png",
                clip={"x": 0, "y": 0, "width": 1280, "height": 720},
            )

            blank = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                io.BytesIO(screenshot),
                left=0, top=0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

        await browser.close()

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
